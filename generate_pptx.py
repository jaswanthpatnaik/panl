import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Color definitions
    BG_COLOR = RGBColor(11, 15, 25)       # Deep Slate/Navy #0B0F19
    CYAN = RGBColor(0, 242, 255)         # primary #00F2FF
    PURPLE = RGBColor(112, 0, 255)       # accent #7000FF
    WHITE = RGBColor(243, 244, 246)      # primary text #F3F4F6
    GRAY = RGBColor(156, 163, 175)       # secondary text #9CA3AF
    
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
    def add_header(slide, title_text):
        # Adds logo and slide title
        tb = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.833), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "PNAL  |  " + title_text
        p.font.name = 'Arial'
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = CYAN
        
    # ==================== SLIDE 1: Title Slide ====================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    
    # Subtitle
    tb_sub = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(0.5))
    p_sub = tb_sub.text_frame.paragraphs[0]
    p_sub.text = "PROACTIVE NETWORK & ANALYSIS LABORATORY"
    p_sub.font.name = 'Arial'
    p_sub.font.size = Pt(14)
    p_sub.font.bold = True
    p_sub.font.color.rgb = CYAN
    p_sub.font.letter_spacing = Pt(3)
    
    # Main Title
    tb_title = slide1.shapes.add_textbox(Inches(1.0), Inches(2.7), Inches(11.333), Inches(2.0))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = "Forensic Malware\nAnalysis Engine"
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(56)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE
    
    # Divider line
    line = slide1.shapes.add_shape(
        1, # Rectangle
        Inches(1.0), Inches(4.8), Inches(2.0), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PURPLE
    line.line.color.rgb = PURPLE
    
    # Description
    tb_desc = slide1.shapes.add_textbox(Inches(1.0), Inches(5.1), Inches(11.333), Inches(1.0))
    p_desc = tb_desc.text_frame.paragraphs[0]
    p_desc.text = "Air-Gapped Document DNA Audit, Active Payload Carving & Zero-Trust Sandbox"
    p_desc.font.name = 'Arial'
    p_desc.font.size = Pt(16)
    p_desc.font.color.rgb = GRAY
    
    # ==================== SLIDE 2: Engine Data Flow ====================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "Unified Analysis Data Pipeline")
    
    # Description Text (Left)
    tb_info2 = slide2.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.0), Inches(5.0))
    tf_info2 = tb_info2.text_frame
    tf_info2.word_wrap = True
    
    p = tf_info2.paragraphs[0]
    p.text = "Ingestion & Analysis Pipeline"
    p.font.name = 'Arial'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(20)
    
    bullets = [
        "Structured Container Parsing: Dissects layouts dynamically, isolating xref and object streams cleanly.",
        "Heuristics & Signatures: Parallel scan loops evaluate custom forensic metadata profiles.",
        "Deep Metadata Extraction: Evaluates author credentials, software versions, and modification trails.",
        "Reputation Correlation: Compares local forensic markers with global intelligence indicators."
    ]
    for b in bullets:
        p2 = tf_info2.add_paragraph()
        p2.text = "• " + b
        p2.font.name = 'Arial'
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
        p2.space_after = Pt(12)
        
    # Flowchart Image (Right)
    img_path2 = "/home/lenovo/.gemini/antigravity/scratch/PNAL_Submission/diagrams/pnal_architecture_flowchart.png"
    if os.path.exists(img_path2):
        slide2.shapes.add_picture(img_path2, Inches(6.25), Inches(1.5), Inches(6.333), Inches(5.0))
        
    # ==================== SLIDE 3: Zero-Trust Sandbox ====================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "Secure Sandbox & Isolated Preview")
    
    # Description Text (Left)
    tb_info3 = slide3.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.0), Inches(5.0))
    tf_info3 = tb_info3.text_frame
    tf_info3.word_wrap = True
    
    p = tf_info3.paragraphs[0]
    p.text = "Digital Blast Shield"
    p.font.name = 'Arial'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(20)
    
    bullets3 = [
        "Cryptographic CSP: Blocks script injections and outbound network connections via native browser policies.",
        "Active Content Stripping: Neutralizes dangerous execution directives (like /JS, /OpenAction) dynamically.",
        "Process Containment: Restricts DOM interactions to keep untrusted elements strictly quarantined.",
        "Analyst View Safety: Ensures absolute document readability without client-side execution risk."
    ]
    for b in bullets3:
        p2 = tf_info3.add_paragraph()
        p2.text = "• " + b
        p2.font.name = 'Arial'
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
        p2.space_after = Pt(12)
        
    # Flowchart Image (Right)
    img_path3 = "/home/lenovo/.gemini/antigravity/scratch/PNAL_Submission/diagrams/pnal_sandbox_flowchart.png"
    if os.path.exists(img_path3):
        slide3.shapes.add_picture(img_path3, Inches(6.25), Inches(1.5), Inches(6.333), Inches(5.0))

    # ==================== SLIDE 4: Forensic Dashboard ====================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Forensic Analyst Hub Dashboard")
    
    # Description Text (Left)
    tb_info4 = slide4.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.0), Inches(5.0))
    tf_info4 = tb_info4.text_frame
    tf_info4.word_wrap = True
    
    p = tf_info4.paragraphs[0]
    p.text = "Enterprise-Grade Triage"
    p.font.name = 'Arial'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(20)
    
    bullets4 = [
        "Suspicion Gauges: Live threat score visualization from local and online checks.",
        "IOC Pivot Grid: Immediate identification of embedded C2 servers, hashes, and URLs.",
        "Behavioral Log: Granular analysis timeline displaying detailed actions step-by-step.",
        "Visual Sandboxing: Real-time isolated rendering side-by-side with structural anomalies."
    ]
    for b in bullets4:
        p2 = tf_info4.add_paragraph()
        p2.text = "• " + b
        p2.font.name = 'Arial'
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
        p2.space_after = Pt(12)
        
    # Mockup Image (Right)
    img_path4 = "/home/lenovo/.gemini/antigravity/scratch/PNAL_Submission/diagrams/pnal_dashboard_mockup.png"
    if os.path.exists(img_path4):
        slide4.shapes.add_picture(img_path4, Inches(6.25), Inches(1.5), Inches(6.333), Inches(5.0))

    # ==================== SLIDE 5: Core Capabilities ====================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Technical Capabilities Breakdown")
    
    # Title
    tb_title5 = slide5.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(0.8))
    p5 = tb_title5.text_frame.paragraphs[0]
    p5.text = "Surgical Threats Investigation"
    p5.font.name = 'Arial'
    p5.font.size = Pt(24)
    p5.font.bold = True
    p5.font.color.rgb = WHITE
    
    # 3 Columns for features
    col_width = Inches(3.644)
    col_gap = Inches(0.45)
    col_top = Inches(2.5)
    col_height = Inches(4.0)
    
    features = [
        {
            "num": "01",
            "title": "Active Payload Carving",
            "desc": "Audits document internal containers dynamically. Exposes and extracts compressed OLE elements, binary streams, scripts, and embedded files to compute SHA256 hashes for triage."
        },
        {
            "num": "02",
            "title": "Forensic Image Audit",
            "desc": "Inspects embedded media blocks to detect web shell triggers and hidden malicious scripts. Discovers steganographic concealment inside high-capacity JPEG, PNG, and SVG vector structures."
        },
        {
            "num": "03",
            "title": "Zero-Day Threat Scans",
            "desc": "Utilizes offline YARA signature compilation and structural container validation. Detects malformed tables and AV-evasion patterns instantly before signatures are indexed globally."
        }
    ]
    
    for i, f in enumerate(features):
        left_pos = Inches(0.75) + i * (col_width + col_gap)
        tb_col = slide5.shapes.add_textbox(left_pos, col_top, col_width, col_height)
        tf_col = tb_col.text_frame
        tf_col.word_wrap = True
        
        # Large Number
        p_num = tf_col.paragraphs[0]
        p_num.text = f["num"]
        p_num.font.name = 'Arial'
        p_num.font.size = Pt(44)
        p_num.font.bold = True
        p_num.font.color.rgb = CYAN
        p_num.space_after = Pt(10)
        
        # Title
        p_title = tf_col.add_paragraph()
        p_title.text = f["title"]
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE
        p_title.space_after = Pt(10)
        
        # Desc
        p_desc = tf_col.add_paragraph()
        p_desc.text = f["desc"]
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = GRAY
        p_desc.line_spacing = 1.3
        
    # ==================== SLIDE 6: Portability & Air-Gapped ====================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "Operational Readiness")
    
    tb_info6 = slide6.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.833), Inches(4.5))
    tf_info6 = tb_info6.text_frame
    tf_info6.word_wrap = True
    
    p = tf_info6.paragraphs[0]
    p.text = "Standalone Windows Compilation"
    p.font.name = 'Arial'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(20)
    
    p2 = tf_info6.add_paragraph()
    p2.text = "PNAL compiles into a fully self-contained, portable Windows binary (pnal.exe) utilizing automated GitHub Actions. All required resources, YARA rules, Flask routes, HTML templates, and static engines are integrated directly inside the executable. This eliminates external runtime setups—providing tactical intelligence response capabilities to forensic specialists on completely air-gapped security operations."
    p2.font.name = 'Arial'
    p2.font.size = Pt(18)
    p2.font.color.rgb = GRAY
    p2.line_spacing = 1.5
    
    # Save Presentation
    dest_path = "/home/lenovo/.gemini/antigravity/scratch/PNAL_Submission/presentation.pptx"
    prs.save(dest_path)
    print(f"[SUCCESS] Presentation saved to {dest_path}")

if __name__ == "__main__":
    create_presentation()
